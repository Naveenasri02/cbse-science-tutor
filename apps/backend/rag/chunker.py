import io
import re
import warnings
import config

# Optional OCR support
try:
    import easyocr
    _ocr_reader = None  # lazy-initialized on first use
    _HAS_OCR = True
except ImportError:
    _HAS_OCR = False


# ── Utilities ────────────────────────────────────────────────────────────


def _estimate_tokens(text: str) -> int:
    """Rough token count: ~4 chars per token."""
    return max(1, len(text) // 4)


def _clean_text(text: str) -> str:
    """Collapse excessive whitespace while preserving structure."""
    text = re.sub(r'\n{3,}', '\n\n', text)
    text = re.sub(r'[ \t]{2,}', ' ', text)
    return text.strip()


def _make_chunk(
    text: str,
    chunk_type: str = "paragraph",
    page: int = 1,
    section: str = "",
    section_hierarchy: list | None = None,
    position: int = 0,
) -> dict:
    """Create a chunk dict with metadata."""
    return {
        "text": text,
        "type": chunk_type,
        "page": page,
        "section": section,
        "section_hierarchy": list(section_hierarchy or []),
        "position": position,
    }


def _recursive_split(text: str, chunk_size: int, chunk_overlap: int) -> list[str]:
    """Split text into chunks at natural boundaries with overlap."""
    separators = ["\n\n", "\n", ". ", "! ", "? ", "; ", ", ", " "]
    chunks: list[str] = []
    start = 0
    text_len = len(text)
    char_chunk = chunk_size * 4  # convert token estimate to chars
    char_overlap = chunk_overlap * 4

    while start < text_len:
        end = min(start + char_chunk, text_len)

        # If not at the end, find a natural break point
        if end < text_len:
            best_break = -1
            for sep in separators:
                search_start = max(start + char_chunk // 2, start)
                idx = text.rfind(sep, search_start, end)
                if idx != -1:
                    best_break = idx + len(sep)
                    break
            if best_break > start:
                end = best_break

        chunk = text[start:end].strip()
        if chunk and _estimate_tokens(chunk) >= 10:
            chunks.append(chunk)

        # Move forward with overlap
        start = max(start + 1, end - char_overlap)

    return chunks


def _table_to_markdown(table_data: list[list]) -> str:
    """Convert table cell data to a markdown table string."""
    if not table_data or not table_data[0]:
        return ""

    rows = [
        [str(cell).strip() if cell is not None else "" for cell in row]
        for row in table_data
    ]
    if not rows:
        return ""

    headers = rows[0]
    num_cols = len(headers)
    col_widths = [max(len(headers[i]), 3) for i in range(num_cols)]
    for row in rows[1:]:
        for i, cell in enumerate(row):
            if i < num_cols:
                col_widths[i] = max(col_widths[i], len(cell))

    lines = []
    # Header
    lines.append(
        "| " + " | ".join(h.ljust(col_widths[i]) for i, h in enumerate(headers)) + " |"
    )
    # Separator
    lines.append(
        "| " + " | ".join("-" * col_widths[i] for i in range(num_cols)) + " |"
    )
    # Data rows
    for row in rows[1:]:
        cells = []
        for i in range(num_cols):
            val = row[i] if i < len(row) else ""
            cells.append(val.ljust(col_widths[i]))
        lines.append("| " + " | ".join(cells) + " |")

    return "\n".join(lines)


# ── PDF: OCR helper ─────────────────────────────────────────────────────


def _ocr_page(page) -> str:
    """Attempt OCR on an image-only page. Returns text or empty string."""
    if not _HAS_OCR:
        warnings.warn(
            "Page appears to be image-only but easyocr is not installed. "
            "Skipping OCR.",
            stacklevel=2,
        )
        return ""
    try:
        global _ocr_reader
        if _ocr_reader is None:
            _ocr_reader = easyocr.Reader(["en"], gpu=False)

        pix = page.get_pixmap(dpi=200)
        img_bytes = pix.tobytes("png")
        results = _ocr_reader.readtext(img_bytes, detail=0)
        return " ".join(results).strip()
    except Exception as exc:
        warnings.warn(f"OCR failed: {exc}", stacklevel=2)
        return ""


# ── PDF: block-level extraction ──────────────────────────────────────────


def _extract_pdf_blocks(file_bytes: bytes) -> list[dict]:
    """Extract structured blocks from a PDF using PyMuPDF dict mode.

    Returns raw block dicts with keys:
        text, type, page, font_size, is_heading
    """
    import fitz

    doc = fitz.open(stream=file_bytes, filetype="pdf")
    blocks: list[dict] = []

    for page_idx in range(len(doc)):
        page = doc[page_idx]
        page_no = page_idx + 1  # 1-indexed

        # ── Tables (extract first so we can skip overlapping text) ──
        table_rects: list = []
        try:
            tables = page.find_tables()
            for table in tables:
                table_data = table.extract()
                if table_data and any(
                    any(cell for cell in row) for row in table_data
                ):
                    md = _table_to_markdown(table_data)
                    if md.strip():
                        blocks.append({
                            "text": md,
                            "type": "table",
                            "page": page_no,
                            "font_size": 0,
                            "is_heading": False,
                        })
                        table_rects.append(fitz.Rect(table.bbox))
        except Exception:
            pass

        # ── Text blocks via dict mode ──
        page_dict = page.get_text("dict", flags=fitz.TEXT_PRESERVE_WHITESPACE)
        text_block_count = 0
        image_block_count = 0

        for block in page_dict.get("blocks", []):
            if block["type"] == 1:  # image block
                image_block_count += 1
                continue

            # Skip text that overlaps an already-extracted table
            block_rect = fitz.Rect(block["bbox"])
            if any(block_rect.intersects(tr) for tr in table_rects):
                continue

            # Collect spans from all lines
            block_text_parts: list[str] = []
            max_font_size = 0.0
            for line in block.get("lines", []):
                line_parts: list[str] = []
                for span in line.get("spans", []):
                    span_text = span.get("text", "")
                    if span_text.strip():
                        line_parts.append(span_text)
                        size = span.get("size", 0)
                        if size > max_font_size:
                            max_font_size = size
                if line_parts:
                    block_text_parts.append("".join(line_parts))

            full_text = "\n".join(block_text_parts).strip()
            if not full_text:
                continue

            text_block_count += 1
            is_heading = max_font_size > 14

            # Classify block type
            if is_heading:
                block_type = "heading"
            elif re.match(
                r'^[\s]*[\u2022\u2023\u25E6\u2043\u2219•\-\*]\s', full_text
            ) or re.match(r'^[\s]*\d+[\.\)]\s', full_text):
                block_type = "list"
            else:
                block_type = "paragraph"

            blocks.append({
                "text": full_text,
                "type": block_type,
                "page": page_no,
                "font_size": max_font_size,
                "is_heading": is_heading,
            })

        # ── OCR fallback for image-only pages ──
        if text_block_count == 0 and image_block_count > 0:
            ocr_text = _ocr_page(page)
            if ocr_text:
                blocks.append({
                    "text": ocr_text,
                    "type": "paragraph",
                    "page": page_no,
                    "font_size": 0,
                    "is_heading": False,
                })

    doc.close()
    return blocks


# ── PDF: section hierarchy ───────────────────────────────────────────────


def _build_section_hierarchy(blocks: list[dict]) -> list[dict]:
    """Annotate blocks with *section*, *section_hierarchy*, and *position*."""
    hierarchy: list[tuple[int, str]] = []  # (level, heading_text)
    current_section = ""
    section_position = 0

    annotated: list[dict] = []
    for block in blocks:
        if block["is_heading"]:
            heading_text = _clean_text(block["text"])
            # level 0 for top-level (>18pt), level 1 for sub-section (>14pt)
            level = 0 if block["font_size"] > 18 else 1

            # Trim hierarchy to current level and append new heading
            hierarchy = [(lv, txt) for lv, txt in hierarchy if lv < level]
            hierarchy.append((level, heading_text))
            current_section = heading_text
            section_position = 0

        block["section"] = current_section
        block["section_hierarchy"] = [txt for _, txt in hierarchy]
        block["position"] = section_position
        section_position += 1
        annotated.append(block)

    return annotated


# ── PDF: structure-aware chunking ────────────────────────────────────────


def _chunk_pdf_blocks(blocks: list[dict]) -> list[dict]:
    """Convert annotated PDF blocks into final chunks.

    Rules:
    - Tables are never split; each is its own chunk (up to 1024 tokens).
    - Headings become their own chunk.
    - Paragraphs / lists are buffered per-section and split at ~CHUNK_SIZE
      tokens with CHUNK_OVERLAP overlap.  Section heading is prepended to
      continuation chunks for context.
    """
    chunks: list[dict] = []
    paragraph_buffer: list[str] = []
    buffer_meta: dict | None = None

    def flush_buffer():
        nonlocal paragraph_buffer, buffer_meta
        if not paragraph_buffer or buffer_meta is None:
            paragraph_buffer = []
            buffer_meta = None
            return

        combined = _clean_text("\n\n".join(paragraph_buffer))

        if _estimate_tokens(combined) <= config.CHUNK_SIZE:
            if _estimate_tokens(combined) >= 10:
                chunks.append(_make_chunk(
                    text=combined,
                    chunk_type=buffer_meta["type"],
                    page=buffer_meta["page"],
                    section=buffer_meta["section"],
                    section_hierarchy=buffer_meta["section_hierarchy"],
                    position=buffer_meta["position"],
                ))
        else:
            section_prefix = buffer_meta["section"]
            sub_chunks = _recursive_split(
                combined, config.CHUNK_SIZE, config.CHUNK_OVERLAP,
            )
            for i, sc in enumerate(sub_chunks):
                # Prepend section heading to overlap chunks for context
                if i > 0 and section_prefix:
                    sc = f"[{section_prefix}]\n{sc}"
                chunks.append(_make_chunk(
                    text=sc,
                    chunk_type=buffer_meta["type"],
                    page=buffer_meta["page"],
                    section=buffer_meta["section"],
                    section_hierarchy=buffer_meta["section_hierarchy"],
                    position=buffer_meta["position"] + i,
                ))

        paragraph_buffer = []
        buffer_meta = None

    for block in blocks:
        if block["type"] == "table":
            flush_buffer()
            table_text = _clean_text(block["text"])
            if _estimate_tokens(table_text) > 1024:
                table_text = table_text[: 1024 * 4] + "\n[… table truncated]"
            if table_text:
                chunks.append(_make_chunk(
                    text=table_text,
                    chunk_type="table",
                    page=block["page"],
                    section=block["section"],
                    section_hierarchy=block["section_hierarchy"],
                    position=block["position"],
                ))

        elif block["type"] == "heading":
            flush_buffer()
            heading_text = _clean_text(block["text"])
            if heading_text and _estimate_tokens(heading_text) >= 3:
                chunks.append(_make_chunk(
                    text=heading_text,
                    chunk_type="heading",
                    page=block["page"],
                    section=block["section"],
                    section_hierarchy=block["section_hierarchy"],
                    position=block["position"],
                ))

        else:  # paragraph or list
            # Flush on section boundary
            if buffer_meta is not None and buffer_meta["section"] != block["section"]:
                flush_buffer()

            if buffer_meta is None:
                buffer_meta = {
                    "type": block["type"],
                    "page": block["page"],
                    "section": block["section"],
                    "section_hierarchy": block["section_hierarchy"],
                    "position": block["position"],
                }

            paragraph_buffer.append(block["text"])

            if _estimate_tokens("\n\n".join(paragraph_buffer)) >= config.CHUNK_SIZE:
                flush_buffer()

    flush_buffer()
    return chunks


def parse_pdf(file_bytes: bytes) -> list[dict]:
    """Extract structured chunks from a PDF with rich metadata."""
    blocks = _extract_pdf_blocks(file_bytes)
    if not blocks:
        raise ValueError("PDF contains no extractable content.")
    annotated = _build_section_hierarchy(blocks)
    return _chunk_pdf_blocks(annotated)


# ── Non-PDF parsers (preserved, return plain text) ───────────────────────


def _parse_docx_xml(file_bytes: bytes) -> str:
    """Fallback: extract text by reading DOCX XML directly from the ZIP."""
    import zipfile
    import xml.etree.ElementTree as ET

    WORD_NS = "{http://schemas.openxmlformats.org/wordprocessingml/2006/main}"
    zf = zipfile.ZipFile(io.BytesIO(file_bytes))

    doc_xml = None
    for name in ["word/document.xml", "word/document2.xml"]:
        if name in zf.namelist():
            doc_xml = zf.read(name)
            break
    if not doc_xml:
        for name in zf.namelist():
            if name.startswith("word/") and name.endswith(".xml") and "document" in name.lower():
                doc_xml = zf.read(name)
                break
    if not doc_xml:
        raise ValueError("Could not find document content in DOCX file")

    root = ET.fromstring(doc_xml)
    paragraphs = []
    for p in root.iter(f"{WORD_NS}p"):
        texts = [t.text for t in p.iter(f"{WORD_NS}t") if t.text]
        if texts:
            paragraphs.append("".join(texts))
    return "\n\n".join(paragraphs)


def parse_docx(file_bytes: bytes) -> str:
    """Extract text from DOCX bytes with multiple fallbacks."""
    # Try python-docx first
    try:
        from docx import Document
        doc = Document(io.BytesIO(file_bytes))
        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
        text = "\n\n".join(paragraphs)
        if text.strip():
            return text
    except Exception:
        pass

    # Fallback: parse DOCX XML directly from ZIP
    try:
        text = _parse_docx_xml(file_bytes)
        if text.strip():
            return text
    except Exception:
        pass

    # Last resort: use PyMuPDF (supports DOCX too)
    try:
        import fitz
        doc = fitz.open(stream=file_bytes, filetype="docx")
        pages = [page.get_text("text") for page in doc]
        doc.close()
        text = "\n\n".join(pages)
        if text.strip():
            return text
    except Exception:
        pass

    raise ValueError("Could not extract text from this Word document. The file may be corrupted.")


def parse_doc(file_bytes: bytes) -> str:
    """Extract text from legacy binary .doc (OLE format)."""
    # Try PyMuPDF with correct filetype
    try:
        import fitz
        doc = fitz.open(stream=file_bytes, filetype="doc")
        pages = [page.get_text("text") for page in doc]
        doc.close()
        text = "\n\n".join(pages)
        if text.strip():
            return text
    except Exception:
        pass

    # Try antiword (Linux CLI tool)
    try:
        import subprocess
        import tempfile
        with tempfile.NamedTemporaryFile(suffix=".doc", delete=False) as tmp:
            tmp.write(file_bytes)
            tmp_path = tmp.name
        result = subprocess.run(
            ["antiword", tmp_path], capture_output=True, text=True, timeout=30
        )
        import os
        os.unlink(tmp_path)
        if result.returncode == 0 and result.stdout.strip():
            return result.stdout
    except Exception:
        pass

    # Try catdoc (another Linux CLI tool)
    try:
        import subprocess
        import tempfile
        with tempfile.NamedTemporaryFile(suffix=".doc", delete=False) as tmp:
            tmp.write(file_bytes)
            tmp_path = tmp.name
        result = subprocess.run(
            ["catdoc", tmp_path], capture_output=True, text=True, timeout=30
        )
        import os
        os.unlink(tmp_path)
        if result.returncode == 0 and result.stdout.strip():
            return result.stdout
    except Exception:
        pass

    # Try olefile raw text extraction
    try:
        import olefile
        ole = olefile.OleFileIO(io.BytesIO(file_bytes))
        if ole.exists("WordDocument"):
            stream = ole.openstream("WordDocument")
            raw = stream.read()
            # Extract printable ASCII/Unicode text segments
            text = raw.decode("utf-8", errors="ignore")
            import re as _re
            # Keep runs of printable chars (3+ length)
            segments = _re.findall(r'[\x20-\x7E\n\r\t]{3,}', text)
            text = " ".join(segments)
            ole.close()
            if len(text.strip()) > 50:
                return text
        ole.close()
    except Exception:
        pass

    raise ValueError(
        "Could not extract text from this .doc file. "
        "Try converting it to .docx or .pdf first."
    )


def parse_txt(file_bytes: bytes) -> str:
    """Extract text from plain text / markdown files."""
    for enc in ("utf-8", "utf-8-sig", "latin-1"):
        try:
            return file_bytes.decode(enc)
        except UnicodeDecodeError:
            continue
    return file_bytes.decode("utf-8", errors="replace")


def parse_csv(file_bytes: bytes) -> str:
    """Convert CSV to readable text (row-per-line)."""
    import csv
    text = parse_txt(file_bytes)
    reader = csv.reader(io.StringIO(text))
    rows = list(reader)
    if not rows:
        return ""
    headers = rows[0]
    lines = []
    for row in rows[1:]:
        parts = [f"{h}: {v}" for h, v in zip(headers, row) if v.strip()]
        if parts:
            lines.append(", ".join(parts))
    return "\n".join(lines) if lines else "\n".join([", ".join(r) for r in rows])


def parse_pptx(file_bytes: bytes) -> str:
    """Extract text from PowerPoint files."""
    from pptx import Presentation
    prs = Presentation(io.BytesIO(file_bytes))
    slides = []
    for i, slide in enumerate(prs.slides, 1):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = para.text.strip()
                    if t:
                        texts.append(t)
        if texts:
            slides.append(f"Slide {i}:\n" + "\n".join(texts))
    return "\n\n".join(slides)


def parse_xlsx(file_bytes: bytes) -> str:
    """Extract text from Excel files."""
    from openpyxl import load_workbook
    wb = load_workbook(io.BytesIO(file_bytes), read_only=True, data_only=True)
    sheets = []
    for ws in wb.worksheets:
        rows = []
        for row in ws.iter_rows(values_only=True):
            cells = [str(c) if c is not None else "" for c in row]
            if any(c.strip() for c in cells):
                rows.append(", ".join(c for c in cells if c.strip()))
        if rows:
            sheets.append(f"Sheet: {ws.title}\n" + "\n".join(rows))
    wb.close()
    return "\n\n".join(sheets)


# ── Non-PDF chunking helper ─────────────────────────────────────────────


def _chunk_plain_text(
    text: str,
    source_type: str = "paragraph",
    page: int = 1,
) -> list[dict]:
    """Chunk plain text and wrap each piece in a metadata dict."""
    text = _clean_text(text)
    if not text:
        return []
    raw_chunks = _recursive_split(text, config.CHUNK_SIZE, config.CHUNK_OVERLAP)
    return [
        _make_chunk(
            text=c, chunk_type=source_type, page=page,
            section="", section_hierarchy=[], position=i,
        )
        for i, c in enumerate(raw_chunks)
    ]


# ── Public API ───────────────────────────────────────────────────────────


SUPPORTED_EXTENSIONS = {
    ".pdf", ".docx", ".doc",
    ".txt", ".md", ".csv",
    ".pptx", ".xlsx",
}


def process_document(file_bytes: bytes, filename: str) -> list[dict]:
    """Parse and chunk a document.

    Returns a list of chunk dicts, each containing::

        {
            "text": str,
            "type": "paragraph" | "table" | "heading" | "list",
            "page": int,          # 1-indexed
            "section": str,
            "section_hierarchy": list[str],
            "position": int,      # 0-indexed within section
        }
    """
    lower = filename.lower()

    # PDF gets the full structure-aware pipeline
    if lower.endswith(".pdf"):
        chunks = parse_pdf(file_bytes)
        if not chunks:
            raise ValueError("Document appears to be empty or contains no extractable text.")
        return chunks

    # All other formats: parse to flat text, then chunk with basic metadata
    if lower.endswith(".docx"):
        text = parse_docx(file_bytes)
    elif lower.endswith(".doc"):
        text = parse_doc(file_bytes)
    elif lower.endswith((".txt", ".md")):
        text = parse_txt(file_bytes)
    elif lower.endswith(".csv"):
        text = parse_csv(file_bytes)
    elif lower.endswith(".pptx"):
        text = parse_pptx(file_bytes)
    elif lower.endswith(".xlsx"):
        text = parse_xlsx(file_bytes)
    else:
        raise ValueError(
            f"Unsupported file format: {filename}. "
            f"Supported: {', '.join(sorted(SUPPORTED_EXTENSIONS))}"
        )

    if not text.strip():
        raise ValueError("Document appears to be empty or contains no extractable text.")

    return _chunk_plain_text(text)
