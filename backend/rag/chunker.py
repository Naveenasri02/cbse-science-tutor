import io
import re
import config


def _estimate_tokens(text: str) -> int:
    """Rough token count: ~4 chars per token."""
    return max(1, len(text) // 4)


def _recursive_split(text: str, chunk_size: int, chunk_overlap: int) -> list[str]:
    """Split text into chunks at natural boundaries with overlap."""
    separators = ["\n\n", "\n", ". ", "! ", "? ", "; ", ", ", " "]
    chunks = []
    start = 0
    text_len = len(text)
    char_chunk = chunk_size * 4  # convert token estimate to chars
    char_overlap = chunk_overlap * 4

    while start < text_len:
        end = min(start + char_chunk, text_len)

        # If not at the end, find a natural break point
        if end < text_len:
            best_break = -1
            for sep in separators:
                # Search for separator near the end of the chunk
                search_start = max(start + char_chunk // 2, start)
                idx = text.rfind(sep, search_start, end)
                if idx != -1:
                    best_break = idx + len(sep)
                    break
            if best_break > start:
                end = best_break

        chunk = text[start:end].strip()
        if chunk and _estimate_tokens(chunk) >= 10:
            chunks.append(chunk)

        # Move forward with overlap
        start = max(start + 1, end - char_overlap)

    return chunks


def parse_pdf(file_bytes: bytes) -> str:
    """Extract text from PDF bytes using PyMuPDF."""
    import fitz  # PyMuPDF
    doc = fitz.open(stream=file_bytes, filetype="pdf")
    pages = []
    for page in doc:
        pages.append(page.get_text("text"))
    doc.close()
    return "\n\n".join(pages)


def _parse_docx_xml(file_bytes: bytes) -> str:
    """Fallback: extract text by reading DOCX XML directly from the ZIP."""
    import zipfile
    import xml.etree.ElementTree as ET

    WORD_NS = "{http://schemas.openxmlformats.org/wordprocessingml/2006/main}"
    zf = zipfile.ZipFile(io.BytesIO(file_bytes))

    doc_xml = None
    for name in ["word/document.xml", "word/document2.xml"]:
        if name in zf.namelist():
            doc_xml = zf.read(name)
            break
    if not doc_xml:
        for name in zf.namelist():
            if name.startswith("word/") and name.endswith(".xml") and "document" in name.lower():
                doc_xml = zf.read(name)
                break
    if not doc_xml:
        raise ValueError("Could not find document content in DOCX file")

    root = ET.fromstring(doc_xml)
    paragraphs = []
    for p in root.iter(f"{WORD_NS}p"):
        texts = [t.text for t in p.iter(f"{WORD_NS}t") if t.text]
        if texts:
            paragraphs.append("".join(texts))
    return "\n\n".join(paragraphs)


def parse_docx(file_bytes: bytes) -> str:
    """Extract text from DOCX bytes with multiple fallbacks."""
    # Try python-docx first
    try:
        from docx import Document
        doc = Document(io.BytesIO(file_bytes))
        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
        text = "\n\n".join(paragraphs)
        if text.strip():
            return text
    except Exception:
        pass

    # Fallback: parse DOCX XML directly from ZIP
    try:
        text = _parse_docx_xml(file_bytes)
        if text.strip():
            return text
    except Exception:
        pass

    # Last resort: use PyMuPDF (supports DOCX too)
    try:
        import fitz
        doc = fitz.open(stream=file_bytes, filetype="docx")
        pages = [page.get_text("text") for page in doc]
        doc.close()
        text = "\n\n".join(pages)
        if text.strip():
            return text
    except Exception:
        pass

    raise ValueError("Could not extract text from this Word document. The file may be corrupted.")


def parse_txt(file_bytes: bytes) -> str:
    """Extract text from plain text / markdown files."""
    for enc in ("utf-8", "utf-8-sig", "latin-1"):
        try:
            return file_bytes.decode(enc)
        except UnicodeDecodeError:
            continue
    return file_bytes.decode("utf-8", errors="replace")


def parse_csv(file_bytes: bytes) -> str:
    """Convert CSV to readable text (row-per-line)."""
    import csv
    text = parse_txt(file_bytes)
    reader = csv.reader(io.StringIO(text))
    rows = list(reader)
    if not rows:
        return ""
    headers = rows[0]
    lines = []
    for row in rows[1:]:
        parts = [f"{h}: {v}" for h, v in zip(headers, row) if v.strip()]
        if parts:
            lines.append(", ".join(parts))
    return "\n".join(lines) if lines else "\n".join([", ".join(r) for r in rows])


def parse_pptx(file_bytes: bytes) -> str:
    """Extract text from PowerPoint files."""
    from pptx import Presentation
    prs = Presentation(io.BytesIO(file_bytes))
    slides = []
    for i, slide in enumerate(prs.slides, 1):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = para.text.strip()
                    if t:
                        texts.append(t)
        if texts:
            slides.append(f"Slide {i}:\n" + "\n".join(texts))
    return "\n\n".join(slides)


def parse_xlsx(file_bytes: bytes) -> str:
    """Extract text from Excel files."""
    from openpyxl import load_workbook
    wb = load_workbook(io.BytesIO(file_bytes), read_only=True, data_only=True)
    sheets = []
    for ws in wb.worksheets:
        rows = []
        for row in ws.iter_rows(values_only=True):
            cells = [str(c) if c is not None else "" for c in row]
            if any(c.strip() for c in cells):
                rows.append(", ".join(c for c in cells if c.strip()))
        if rows:
            sheets.append(f"Sheet: {ws.title}\n" + "\n".join(rows))
    wb.close()
    return "\n\n".join(sheets)


SUPPORTED_EXTENSIONS = {
    ".pdf", ".docx", ".doc",
    ".txt", ".md", ".csv",
    ".pptx", ".xlsx",
}


def parse_document(file_bytes: bytes, filename: str) -> str:
    """Auto-detect format and extract text."""
    lower = filename.lower()
    if lower.endswith(".pdf"):
        return parse_pdf(file_bytes)
    elif lower.endswith((".docx", ".doc")):
        return parse_docx(file_bytes)
    elif lower.endswith((".txt", ".md")):
        return parse_txt(file_bytes)
    elif lower.endswith(".csv"):
        return parse_csv(file_bytes)
    elif lower.endswith(".pptx"):
        return parse_pptx(file_bytes)
    elif lower.endswith(".xlsx"):
        return parse_xlsx(file_bytes)
    else:
        raise ValueError(f"Unsupported file format: {filename}. Supported: {', '.join(sorted(SUPPORTED_EXTENSIONS))}")


def chunk_text(text: str) -> list[str]:
    """Split document text into overlapping chunks."""
    # Clean up excessive whitespace
    text = re.sub(r'\n{3,}', '\n\n', text)
    text = re.sub(r' {2,}', ' ', text)

    return _recursive_split(text, config.CHUNK_SIZE, config.CHUNK_OVERLAP)


def process_document(file_bytes: bytes, filename: str) -> list[str]:
    """Parse and chunk a document. Returns list of text chunks."""
    text = parse_document(file_bytes, filename)
    if not text.strip():
        raise ValueError("Document appears to be empty or contains no extractable text.")
    return chunk_text(text)
